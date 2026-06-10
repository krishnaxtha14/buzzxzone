"""
BuzzXZone Presentation Generator
13-slide visual presentation with forest/neon theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Brand colours ──────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0A, 0x14, 0x0A)   # near-black green
NEON_GRN  = RGBColor(0x00, 0xFF, 0x88)   # neon green accent
NEON_BLU  = RGBColor(0x00, 0xD4, 0xFF)   # neon cyan accent
NEON_PRP  = RGBColor(0xBF, 0x5F, 0xFF)   # neon purple
CARD_BG   = RGBColor(0x0D, 0x2B, 0x1A)   # dark green card
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW    = RGBColor(0xFF, 0xD7, 0x00)
ORANGE    = RGBColor(0xFF, 0x7A, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper utilities ────────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, w, h, fill_color=None, border_color=None,
        border_width=Pt(0), radius=None):
    """Add a rounded/plain rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(w), Inches(h)
    )
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, w, h,
        font_size=Pt(18), color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def icon_box(slide, emoji, label, left, top, w=1.8, h=1.6,
             fill=CARD_BG, accent=NEON_GRN):
    """A small icon card with emoji + label."""
    box(slide, left, top, w, h, fill_color=fill, border_color=accent, border_width=Pt(1.5))
    txt(slide, emoji, left, top + 0.1, w, 0.9,
        font_size=Pt(32), align=PP_ALIGN.CENTER)
    txt(slide, label, left, top + 0.9, w, 0.55,
        font_size=Pt(13), color=accent, bold=True, align=PP_ALIGN.CENTER)


def badge(slide, text, left, top, w=1.5, h=0.42,
          bg=NEON_GRN, fg=BG_DARK):
    """Small filled badge/pill."""
    box(slide, left, top, w, h, fill_color=bg, border_color=None)
    txt(slide, text, left, top + 0.03, w, 0.38,
        font_size=Pt(11), color=fg, bold=True, align=PP_ALIGN.CENTER)


def section_header(slide, title, subtitle=""):
    """Top-left section title block."""
    txt(slide, title, 0.4, 0.2, 12, 0.6,
        font_size=Pt(30), color=NEON_GRN, bold=True)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.75, 10, 0.45,
            font_size=Pt(16), color=RGBColor(0xAA, 0xFF, 0xCC), italic=True)
    # decorative accent line
    box(slide, 0.4, 1.18, 4.5, 0.04, fill_color=NEON_GRN)


def timeline_node(slide, step, label, desc, left, top,
                  accent=NEON_GRN, num_color=BG_DARK):
    """Vertical timeline node."""
    # circle number
    circle = slide.shapes.add_shape(9, Inches(left), Inches(top), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    p = circle.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(step)
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = num_color
    # label + desc
    txt(slide, label, left + 0.65, top - 0.05, 3.8, 0.38,
        font_size=Pt(14), color=WHITE, bold=True)
    txt(slide, desc, left + 0.65, top + 0.3, 3.8, 0.38,
        font_size=Pt(11), color=RGBColor(0xCC, 0xFF, 0xDD))


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)

# big neon glow box behind title
box(s, 1.5, 1.6, 10.3, 2.4, fill_color=RGBColor(0x00, 0x40, 0x20), border_color=NEON_GRN, border_width=Pt(2))

txt(s, "🌿  buzzXzone", 1.7, 1.75, 10, 1.0,
    font_size=Pt(56), color=NEON_GRN, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Eco Learning Hub", 1.7, 2.7, 10, 0.6,
    font_size=Pt(26), color=WHITE, align=PP_ALIGN.CENTER)

# tag line badges
badge(s, "🐍 Snake Quiz",    2.0, 4.2, 2.0, 0.5, bg=RGBColor(0x00, 0x80, 0x40), fg=WHITE)
badge(s, "🧮 Math Quiz",     4.3, 4.2, 2.0, 0.5, bg=RGBColor(0x00, 0x60, 0xA0), fg=WHITE)
badge(s, "🔒 Cyber Quiz",    6.6, 4.2, 2.0, 0.5, bg=RGBColor(0x60, 0x00, 0xA0), fg=WHITE)
badge(s, "🃏 Memory Match",  8.9, 4.2, 2.2, 0.5, bg=RGBColor(0xA0, 0x40, 0x00), fg=WHITE)

txt(s, "An educational gamified web platform built with Flask + PostgreSQL",
    1.5, 5.1, 10.3, 0.5, font_size=Pt(15), color=RGBColor(0x88, 0xCC, 0xAA),
    align=PP_ALIGN.CENTER, italic=True)

# bottom link
txt(s, "🔗  buzzxzone.vercel.app", 4.5, 6.5, 5, 0.4,
    font_size=Pt(14), color=NEON_BLU, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Project Overview
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "What is buzzXzone?", "Gamified eco-education in your browser")

# three big feature cards
cards = [
    ("🌱", "Eco Education",  "Learn ecology &\ncybersecurity\nthrough play"),
    ("🎮", "Game Engine",     "Snake, Quiz, Memory\ngames — all in one\nplatform"),
    ("🏆", "Gamification",   "Points, unlockables\n& high-score\nleaderboard"),
]
for i, (em, title, desc) in enumerate(cards):
    lft = 0.6 + i * 4.2
    box(s, lft, 1.5, 3.9, 4.8, fill_color=CARD_BG, border_color=NEON_GRN, border_width=Pt(1.5))
    txt(s, em, lft, 1.7, 3.9, 1.1, font_size=Pt(52), align=PP_ALIGN.CENTER)
    txt(s, title, lft, 2.85, 3.9, 0.55, font_size=Pt(20), color=NEON_GRN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  lft, 3.45, 3.9, 1.8, font_size=Pt(15), color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "🌐  buzzxzone.vercel.app", 4.5, 6.6, 5, 0.4,
    font_size=Pt(13), color=NEON_BLU, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "Tech Stack", "Full-stack web application")

stack = [
    ("🐍", "Python / Flask",   "Backend logic, routing,\nOTP email, session auth",   NEON_GRN),
    ("🐘", "PostgreSQL",        "Supabase cloud DB\nusers + scores",                  NEON_BLU),
    ("🌐", "HTML / CSS / JS",   "Canvas 2D games,\nneon forest UI theme",             NEON_PRP),
    ("☁️",  "Vercel",            "Serverless deployment\nCI/CD auto-deploy",            YELLOW),
]

for i, (em, title, desc, color) in enumerate(stack):
    row = i // 2
    col = i % 2
    lft = 0.6 + col * 6.5
    top = 1.45 + row * 2.7
    box(s, lft, top, 6.1, 2.4, fill_color=CARD_BG, border_color=color, border_width=Pt(1.5))
    txt(s, em,    lft + 0.2, top + 0.25, 1.2, 1.1, font_size=Pt(44), align=PP_ALIGN.CENTER)
    txt(s, title, lft + 1.4, top + 0.25, 4.5, 0.55, font_size=Pt(20), color=color, bold=True)
    txt(s, desc,  lft + 1.4, top + 0.85, 4.5, 1.0,  font_size=Pt(14), color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Authentication Flow
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "Authentication", "Secure OTP email verification")

steps = [
    ("✉️",  "Register",       "Enter username,\nemail & password"),
    ("📧", "OTP Sent",        "6-digit code via\nGmail SMTP"),
    ("🔑", "Verify",          "Enter OTP to\nactivate account"),
    ("🏠", "Dashboard",       "Access all games\n& your scores"),
]

for i, (em, label, desc) in enumerate(steps):
    lft = 0.5 + i * 3.2
    box(s, lft, 1.5, 2.9, 3.8, fill_color=CARD_BG, border_color=NEON_BLU, border_width=Pt(1.5))
    txt(s, em,    lft, 1.65, 2.9, 1.0, font_size=Pt(44), align=PP_ALIGN.CENTER)
    txt(s, label, lft, 2.7,  2.9, 0.55, font_size=Pt(18), color=NEON_BLU, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  lft, 3.3,  2.9, 1.4,  font_size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
    # arrow (except last)
    if i < 3:
        txt(s, "→", lft + 2.9, 2.85, 0.35, 0.5, font_size=Pt(24), color=NEON_BLU, align=PP_ALIGN.CENTER)

# password reset
box(s, 0.5, 5.6, 12.3, 0.9, fill_color=RGBColor(0x0A, 0x28, 0x3A), border_color=NEON_BLU, border_width=Pt(1))
txt(s, "🔒  Forgot Password?  →  Email OTP  →  Verify  →  Set New Password",
    0.7, 5.7, 12, 0.6, font_size=Pt(15), color=NEON_BLU, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "Dashboard", "Game hub with animated capybara mascot")

# mockup "screen"
box(s, 0.5, 1.45, 6.5, 5.6, fill_color=RGBColor(0x0D, 0x1F, 0x10),
    border_color=NEON_GRN, border_width=Pt(2))

# browser bar
box(s, 0.5, 1.45, 6.5, 0.4, fill_color=RGBColor(0x15, 0x35, 0x1F))
txt(s, "🔗 buzzxzone.vercel.app/dashboard", 0.7, 1.5, 6, 0.32,
    font_size=Pt(10), color=NEON_BLU)

txt(s, "🌿  Welcome, Player!", 0.7, 2.0, 6, 0.5,
    font_size=Pt(20), color=NEON_GRN, bold=True, align=PP_ALIGN.CENTER)

# capybara area
box(s, 0.7, 2.55, 6, 1.3, fill_color=RGBColor(0x05, 0x28, 0x14), border_color=NEON_GRN, border_width=Pt(1))
txt(s, "🦫  Capy says:", 0.9, 2.65, 3, 0.35, font_size=Pt(13), color=NEON_GRN, bold=True)
txt(s, '"Keep learning, every point counts!"', 0.9, 3.0, 5.6, 0.5,
    font_size=Pt(12), color=WHITE, italic=True)

# game cards row
for j, (em, nm, clr) in enumerate([("🐍","Snake",NEON_GRN),("🧮","Math",NEON_BLU),("🔒","Cyber",NEON_PRP)]):
    lft = 0.7 + j * 2.0
    box(s, lft, 4.0, 1.7, 1.6, fill_color=CARD_BG, border_color=clr, border_width=Pt(1.2))
    txt(s, em, lft, 4.1, 1.7, 0.8, font_size=Pt(28), align=PP_ALIGN.CENTER)
    txt(s, nm, lft, 4.9, 1.7, 0.4, font_size=Pt(11), color=clr, bold=True, align=PP_ALIGN.CENTER)

# progress bar
box(s, 0.7, 5.75, 6.0, 0.3, fill_color=RGBColor(0x05, 0x20, 0x10))
box(s, 0.7, 5.75, 3.0, 0.3, fill_color=NEON_GRN)
txt(s, "Score: 50 / 100 pts  →  Unlock Memory Match", 0.7, 6.1, 6, 0.3,
    font_size=Pt(10), color=WHITE)

# feature callouts
features = [
    ("🦫", "Capybara Mascot",    "10 rotating\nmessages on canvas"),
    ("✨", "Neon Mouse Trail",    "Particle cursor\neffect on all pages"),
    ("📊", "Progress Bar",       "Tracks score toward\nMemory Match unlock"),
    ("🌿", "Forest UI Theme",    "Dark-green animated\ncards & float effects"),
]
for i, (em, label, desc) in enumerate(features):
    lft = 7.3 + (i % 2) * 3.0
    top = 1.5 + (i // 2) * 2.7
    box(s, lft, top, 2.7, 2.3, fill_color=CARD_BG, border_color=NEON_GRN, border_width=Pt(1.2))
    txt(s, em,    lft, top + 0.1,  2.7, 0.8, font_size=Pt(34), align=PP_ALIGN.CENTER)
    txt(s, label, lft, top + 0.9,  2.7, 0.45, font_size=Pt(13), color=NEON_GRN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  lft, top + 1.35, 2.7, 0.8,  font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Eco Snake Quiz
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "🐍  Eco Snake Quiz", "Auto-playing snake controlled by your answers")

# snake game mockup
box(s, 0.5, 1.45, 6.3, 5.6, fill_color=RGBColor(0x02, 0x12, 0x06),
    border_color=NEON_GRN, border_width=Pt(2))
txt(s, "🐍  SNAKE QUIZ", 0.5, 1.5, 6.3, 0.5,
    font_size=Pt(18), color=NEON_GRN, bold=True, align=PP_ALIGN.CENTER)

# snake grid
grid_em = [
    "🟩🟩🟩⬛⬛",
    "⬛⬛🟩⬛⬛",
    "⬛⬛🟩🟩🟩",
    "⬛⬛⬛⬛🍎",
]
for r, row in enumerate(grid_em):
    txt(s, row, 1.0, 2.15 + r * 0.55, 5.5, 0.5, font_size=Pt(22), align=PP_ALIGN.CENTER)

# question box
box(s, 0.7, 4.3, 5.8, 0.8, fill_color=RGBColor(0x0A, 0x28, 0x0A), border_color=NEON_GRN, border_width=Pt(1))
txt(s, '❓  "Which gas do plants absorb?"', 0.9, 4.38, 5.5, 0.6,
    font_size=Pt(13), color=WHITE, italic=True)

# answers
for j, (opt, correct) in enumerate([("CO₂ ✓","A"), ("O₂","B"), ("N₂","C"), ("H₂O","D")]):
    clr = NEON_GRN if "✓" in opt else RGBColor(0x44, 0x88, 0x55)
    bg  = RGBColor(0x05, 0x35, 0x12) if "✓" in opt else CARD_BG
    box(s, 0.7 + (j%2)*2.9, 5.2 + (j//2)*0.55, 2.6, 0.45, fill_color=bg, border_color=clr, border_width=Pt(1))
    txt(s, f"{correct}. {opt}", 0.85 + (j%2)*2.9, 5.25 + (j//2)*0.55, 2.4, 0.38,
        font_size=Pt(12), color=clr, bold="✓" in opt)

# callouts
rules = [
    ("✅ Correct",   "+10 pts\nSnake GROWS",  NEON_GRN),
    ("❌ Wrong",      "-pts\nSnake SHRINKS",  RGBColor(0xFF,0x44,0x44)),
    ("⏱ Timer",     "20 seconds\nper question", YELLOW),
    ("🔀 Shuffle",  "Questions\nrandomised",   NEON_BLU),
]
for i, (label, desc, clr) in enumerate(rules):
    lft = 7.1 + (i%2)*3.1
    top = 1.5 + (i//2)*2.7
    box(s, lft, top, 2.85, 2.3, fill_color=CARD_BG, border_color=clr, border_width=Pt(1.5))
    txt(s, label, lft, top + 0.1,  2.85, 0.5,  font_size=Pt(16), color=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  lft, top + 0.65, 2.85, 1.5,  font_size=Pt(22), color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Math Quiz
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "🧮  Math Quiz", "Three difficulty tiers with 10 pts per correct answer")

# difficulty cards
levels = [
    ("😊", "Easy",    "Basic arithmetic\n+ / - / ×",          NEON_GRN,  "10 Questions"),
    ("😤", "Medium",  "Algebra &\npercentages",                 NEON_BLU,  "10 Questions"),
    ("🔥", "Hard",    "Complex equations\n& problem solving",   ORANGE,    "10 Questions"),
]
for i, (em, lvl, desc, clr, count) in enumerate(levels):
    lft = 0.6 + i * 4.2
    box(s, lft, 1.45, 3.9, 4.5, fill_color=CARD_BG, border_color=clr, border_width=Pt(2))
    txt(s, em,    lft, 1.6,  3.9, 1.0, font_size=Pt(52), align=PP_ALIGN.CENTER)
    txt(s, lvl,   lft, 2.65, 3.9, 0.55, font_size=Pt(22), color=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  lft, 3.25, 3.9, 1.2,  font_size=Pt(15), color=WHITE, align=PP_ALIGN.CENTER)
    badge(s, count, lft + 0.95, 4.55, 2.0, 0.42, bg=clr, fg=BG_DARK)

# shared quiz engine note
box(s, 0.6, 6.1, 12.1, 0.8, fill_color=RGBColor(0x0A, 0x22, 0x36), border_color=NEON_BLU, border_width=Pt(1))
txt(s, "⚙️  Powered by the shared quiz.js engine  •  Questions shuffled every session  •  20s timer per question",
    0.8, 6.22, 11.8, 0.55, font_size=Pt(14), color=NEON_BLU, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Cyber Security Quiz
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "🔒  Eco Cyber-Security Quiz", "Digital safety framed around eco habits")

# left: topic areas
topics = [
    "🌐  Safe browsing & phishing awareness",
    "🔑  Password hygiene & 2FA",
    "📧  Email scams & social engineering",
    "🛡️   Data privacy & eco-digital habits",
    "🔍  Malware & suspicious links",
]
box(s, 0.4, 1.45, 5.5, 5.5, fill_color=CARD_BG, border_color=NEON_PRP, border_width=Pt(1.5))
txt(s, "Topics Covered", 0.6, 1.55, 5, 0.45, font_size=Pt(16), color=NEON_PRP, bold=True)
for i, t in enumerate(topics):
    txt(s, t, 0.6, 2.1 + i * 0.85, 5.0, 0.65, font_size=Pt(14), color=WHITE)

# right: difficulty + stats
for i, (lvl, clr, em) in enumerate([("Easy", NEON_GRN,"😊"),("Medium",NEON_BLU,"😤"),("Hard",ORANGE,"🔥")]):
    lft = 6.3 + (i % 2) * 3.4
    top = 1.5 + (i // 2) * 2.5
    box(s, lft, top, 3.0, 2.0, fill_color=CARD_BG, border_color=clr, border_width=Pt(1.5))
    txt(s, em,  lft, top + 0.05, 3.0, 0.9, font_size=Pt(36), align=PP_ALIGN.CENTER)
    txt(s, lvl, lft, top + 1.0,  3.0, 0.55, font_size=Pt(18), color=clr, bold=True, align=PP_ALIGN.CENTER)

# stat strip
box(s, 6.3, 6.15, 6.5, 0.9, fill_color=RGBColor(0x14, 0x0A, 0x28), border_color=NEON_PRP, border_width=Pt(1))
txt(s, "10 pts per answer  •  20s timer  •  Randomised each session",
    6.5, 6.3, 6.1, 0.55, font_size=Pt(13), color=NEON_PRP, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — Memory Match
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "🃏  Memory Match", "Unlock at 100 pts — 4×4 card flip game")

# left: game grid mockup
box(s, 0.5, 1.45, 6.3, 5.6, fill_color=RGBColor(0x08, 0x14, 0x22),
    border_color=NEON_BLU, border_width=Pt(2))
txt(s, "4 × 4 GRID", 0.5, 1.5, 6.3, 0.5,
    font_size=Pt(16), color=NEON_BLU, bold=True, align=PP_ALIGN.CENTER)

card_icons = ["🌊","🌊","🌿","🌿","🦋","🦋","🌸","🌸",
              "🌍","🌍","🌞","🌞","♻️","♻️","🐝","🐝"]
for row in range(4):
    for col in range(4):
        idx = row * 4 + col
        lft = 0.65 + col * 1.45
        top = 2.1 + row * 1.1
        is_flipped = idx in [0, 1, 4, 5]
        bg = RGBColor(0x00, 0x55, 0x99) if is_flipped else RGBColor(0x05, 0x28, 0x40)
        bc = NEON_BLU if is_flipped else RGBColor(0x22, 0x66, 0x88)
        box(s, lft, top, 1.25, 0.9, fill_color=bg, border_color=bc, border_width=Pt(1))
        em_txt = card_icons[idx] if is_flipped else "❓"
        txt(s, em_txt, lft, top + 0.05, 1.25, 0.75, font_size=Pt(24), align=PP_ALIGN.CENTER)

# right: how-to
rules = [
    ("🔒", "Locked",       "Requires 100 pts\nto unlock"),
    ("🔄", "Flip Cards",   "Find matching\npairs of icons"),
    ("✅", "Match Found",  "Cards stay revealed\n+10 pts"),
    ("🏁", "Win Condition","Match all 8 pairs\nto complete"),
]
for i, (em, label, desc) in enumerate(rules):
    lft = 7.1 + (i % 2) * 3.1
    top = 1.5 + (i // 2) * 2.7
    box(s, lft, top, 2.85, 2.3, fill_color=CARD_BG, border_color=NEON_BLU, border_width=Pt(1.5))
    txt(s, em,    lft, top + 0.1,  2.85, 0.8,  font_size=Pt(34), align=PP_ALIGN.CENTER)
    txt(s, label, lft, top + 0.9,  2.85, 0.45, font_size=Pt(14), color=NEON_BLU, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  lft, top + 1.35, 2.85, 0.8,  font_size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — Gamification & Scoring
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "🏆  Gamification", "Points, unlocks & high-score tracking")

# big trophy + score bar
txt(s, "🏆", 0.5, 1.4, 2.5, 1.5, font_size=Pt(80), align=PP_ALIGN.CENTER)

# score progress bar
box(s, 3.0, 1.55, 9.8, 0.65, fill_color=RGBColor(0x05, 0x20, 0x10), border_color=NEON_GRN, border_width=Pt(1))
box(s, 3.0, 1.55, 4.9, 0.65, fill_color=NEON_GRN)
txt(s, "50 pts", 3.1, 1.6, 1.5, 0.5, font_size=Pt(12), color=BG_DARK, bold=True)
txt(s, "🔒 Memory Match unlocks at 100 pts", 7.6, 1.6, 5, 0.5, font_size=Pt(12), color=WHITE)

# milestone boxes
milestones = [
    ("10 pts",  "One correct\nanswer",       "🎯", NEON_GRN),
    ("50 pts",  "Halfway\nto unlock",        "⚡", NEON_BLU),
    ("100 pts", "Memory Match\nUNLOCKED",    "🃏", YELLOW),
    ("∞ pts",   "High Score\nPersists",      "💾", NEON_PRP),
]
for i, (score, desc, em, clr) in enumerate(milestones):
    lft = 0.5 + i * 3.2
    box(s, lft, 2.5, 3.0, 3.5, fill_color=CARD_BG, border_color=clr, border_width=Pt(1.5))
    txt(s, em,    lft, 2.6,  3.0, 1.0,  font_size=Pt(40), align=PP_ALIGN.CENTER)
    txt(s, score, lft, 3.65, 3.0, 0.55, font_size=Pt(22), color=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  lft, 4.25, 3.0, 1.0,  font_size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

# note
txt(s, "⚙️  Scores stored in PostgreSQL  •  High score never decreases  •  memory_unlocked flag auto-set",
    0.5, 6.25, 12.3, 0.5, font_size=Pt(12), color=RGBColor(0x88,0xCC,0xAA),
    align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — Database & Backend Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "⚙️  Backend Architecture", "Flask + PostgreSQL (Supabase) + Vercel")

# flow diagram boxes
components = [
    (0.4,  2.5, "🌐\nBrowser",      NEON_GRN,  RGBColor(0x05,0x28,0x14)),
    (3.2,  2.5, "🐍\nFlask App",    NEON_BLU,  RGBColor(0x05,0x14,0x28)),
    (6.0,  2.5, "🐘\nPostgreSQL",   NEON_PRP,  RGBColor(0x14,0x05,0x28)),
    (8.8,  2.5, "📧\nGmail SMTP",   YELLOW,    RGBColor(0x28,0x20,0x00)),
]
for lft, top, label, clr, bg in components:
    box(s, lft, top, 2.4, 2.0, fill_color=bg, border_color=clr, border_width=Pt(1.8))
    txt(s, label, lft, top + 0.2, 2.4, 1.6, font_size=Pt(22), color=clr,
        bold=True, align=PP_ALIGN.CENTER)
    if lft < 8.8:
        txt(s, "→", lft + 2.4, top + 0.7, 0.4, 0.6, font_size=Pt(24), color=WHITE, align=PP_ALIGN.CENTER)

# DB schema box
box(s, 0.4, 4.8, 5.4, 2.3, fill_color=RGBColor(0x14,0x05,0x28), border_color=NEON_PRP, border_width=Pt(1.5))
txt(s, "users table", 0.6, 4.88, 5.0, 0.45, font_size=Pt(15), color=NEON_PRP, bold=True)
schema_fields = [
    "id  •  username  •  email",
    "password (hashed)  •  high_score",
    "memory_unlocked (boolean)",
]
for i, f in enumerate(schema_fields):
    txt(s, f, 0.6, 5.38 + i * 0.5, 5.0, 0.45, font_size=Pt(12), color=WHITE)

# routes box
box(s, 6.2, 4.8, 6.7, 2.3, fill_color=RGBColor(0x05,0x14,0x28), border_color=NEON_BLU, border_width=Pt(1.5))
txt(s, "Key Routes", 6.4, 4.88, 6.3, 0.45, font_size=Pt(15), color=NEON_BLU, bold=True)
routes = ["/login  /register  /verify","/games/<type>/<difficulty>","/games/memory  /health"]
for i, r in enumerate(routes):
    txt(s, r, 6.4, 5.38 + i * 0.5, 6.3, 0.45, font_size=Pt(12), color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — Deployment & Project Structure
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "☁️  Deployment", "Vercel Serverless  •  Auto CI/CD from GitHub")

# left: file tree
box(s, 0.4, 1.45, 5.5, 5.6, fill_color=CARD_BG, border_color=NEON_GRN, border_width=Pt(1.5))
txt(s, "📁  Project Structure", 0.6, 1.55, 5, 0.45, font_size=Pt(15), color=NEON_GRN, bold=True)
tree = [
    ("📄", "app.py",              "Flask routes + DB logic"),
    ("⚙️",  "vercel.json",         "Serverless config"),
    ("📋", "requirements.txt",    "Python dependencies"),
    ("🗂️",  "questions/*.json",    "7 question banks"),
    ("🎨", "static/style.css",    "Forest neon theme"),
    ("🎮", "static/game.js",      "Snake game engine"),
    ("📝", "templates/*.html",    "13 Jinja2 templates"),
]
for i, (em, name, desc) in enumerate(tree):
    txt(s, f"{em}  {name}", 0.7, 2.15 + i * 0.68, 3.0, 0.5, font_size=Pt(12), color=NEON_GRN, bold=True)
    txt(s, desc,            3.7, 2.15 + i * 0.68, 2.0, 0.5, font_size=Pt(11), color=WHITE)

# right: deploy flow
deploy_steps = [
    ("💻", "Push to\nGitHub",     "git push origin main"),
    ("⚙️",  "Vercel\nCI/CD",       "Auto-build triggers"),
    ("🌐", "Live on\nVercel",     "buzzxzone.vercel.app"),
    ("🔗", "Supabase\nDB",        "PostgreSQL pooler\nport 6543"),
]
for i, (em, label, note) in enumerate(deploy_steps):
    lft = 6.3 + (i % 2) * 3.4
    top = 1.5 + (i // 2) * 2.8
    box(s, lft, top, 3.0, 2.3, fill_color=CARD_BG, border_color=NEON_BLU, border_width=Pt(1.5))
    txt(s, em,    lft, top + 0.05, 3.0, 0.85, font_size=Pt(36), align=PP_ALIGN.CENTER)
    txt(s, label, lft, top + 0.9,  3.0, 0.55, font_size=Pt(14), color=NEON_BLU, bold=True, align=PP_ALIGN.CENTER)
    txt(s, note,  lft, top + 1.45, 3.0, 0.65, font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

# env vars note
box(s, 0.4, 7.05, 12.5, 0.45, fill_color=RGBColor(0x10,0x20,0x10), border_color=NEON_GRN, border_width=Pt(1))
txt(s, "🔐  Env vars: DATABASE_URL  •  FLASK_SECRET_KEY  •  GMAIL_USER  •  GMAIL_PASSWORD",
    0.6, 7.1, 12.1, 0.35, font_size=Pt(12), color=NEON_GRN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)

# large glowing box
box(s, 1.0, 1.2, 11.3, 4.5, fill_color=RGBColor(0x00, 0x30, 0x18),
    border_color=NEON_GRN, border_width=Pt(3))

txt(s, "🌿  Thank You!", 1.0, 1.5, 11.3, 1.4,
    font_size=Pt(66), color=NEON_GRN, bold=True, align=PP_ALIGN.CENTER)

txt(s, "It's been a journey building eco-education through play.",
    1.2, 2.95, 11.0, 0.65, font_size=Pt(20), color=WHITE,
    align=PP_ALIGN.CENTER, italic=True)

txt(s, "Questions? Feedback? Let's talk  🎮",
    1.2, 3.6, 11.0, 0.55, font_size=Pt(18), color=RGBColor(0xAA,0xFF,0xCC),
    align=PP_ALIGN.CENTER)

# bottom badges
badge(s, "🐍 Snake",     1.3,  5.9, 1.9, 0.5, bg=RGBColor(0x00,0x80,0x40), fg=WHITE)
badge(s, "🧮 Math",      3.4,  5.9, 1.9, 0.5, bg=RGBColor(0x00,0x60,0xA0), fg=WHITE)
badge(s, "🔒 Cyber",     5.5,  5.9, 1.9, 0.5, bg=RGBColor(0x60,0x00,0xA0), fg=WHITE)
badge(s, "🃏 Memory",    7.6,  5.9, 1.9, 0.5, bg=RGBColor(0xA0,0x40,0x00), fg=WHITE)
badge(s, "☁️  Deployed", 9.7,  5.9, 2.0, 0.5, bg=RGBColor(0x00,0x50,0x70), fg=WHITE)

txt(s, "🔗  buzzxzone.vercel.app", 3.5, 6.8, 6.5, 0.5,
    font_size=Pt(18), color=NEON_BLU, bold=True, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────────
out = r"c:\Users\sande\OneDrive\Documents\GitHub\buzzxzone\buzzxzone_presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
